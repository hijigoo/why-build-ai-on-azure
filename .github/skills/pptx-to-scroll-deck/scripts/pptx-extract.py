#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx-extract.py — .pptx 원본에서 슬라이드 콘텐츠를 구조화해 뽑아낸다.

    python3 pptx-extract.py 원본.pptx [--img-dir images] [--json out.json]

슬라이드마다 다음을 추출한다.
  - 제목(제목 placeholder 우선, 없으면 첫 텍스트)
  - 본문 텍스트(들여쓰기 수준 유지 — 카드/리스트/표로 옮길 재료)
  - 표(있으면 행렬 그대로)
  - 발표자 노트(→ scroll-deck 의 설명란 재료)
  - 삽입 이미지(→ img-dir 로 저장, <img src> 로 참조)

이 스크립트는 '재료'를 뽑을 뿐이다. 슬라이드 조각(chapter)으로 옮기는 편집·요약·
사실 확인은 SKILL.md 의 Phase 0~2 를 따른다. 특히 GA/Preview·제품명은 원본을
그대로 믿지 말고 공식 문서로 재확인한다.

의존성: python-pptx (pip install python-pptx). 빠른 텍스트만 필요하면
markitdown(python -m markitdown 원본.pptx)으로 대체할 수 있으나, 표/이미지/노트
구조가 필요하면 이 스크립트를 쓴다.
"""
import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx 가 필요합니다:  pip install python-pptx")


def shape_text_lines(shape):
    """텍스트 프레임을 들여쓰기 수준과 함께 줄 리스트로."""
    lines = []
    if not shape.has_text_frame:
        return lines
    for para in shape.text_frame.paragraphs:
        txt = "".join(r.text for r in para.runs).strip()
        if not txt:
            continue
        lvl = para.level or 0
        lines.append({"level": lvl, "text": txt})
    return lines


def extract_table(shape):
    tbl = shape.table
    rows = []
    for r in tbl.rows:
        rows.append([c.text.strip() for c in r.cells])
    return rows


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", help="원본 .pptx 경로")
    ap.add_argument("--img-dir", default="images", help="추출 이미지 저장 폴더")
    ap.add_argument("--json", default=None, help="구조화 결과를 JSON 으로도 저장")
    args = ap.parse_args()

    src = Path(args.pptx)
    if not src.exists():
        sys.exit("파일이 없습니다: %s" % src)

    prs = Presentation(str(src))
    img_dir = Path(args.img_dir)
    slides_out = []
    img_count = 0

    for si, slide in enumerate(prs.slides, start=1):
        title = ""
        body = []
        tables = []
        images = []

        # 제목 placeholder 우선
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            # 이미지
            if getattr(shape, "image", None) is not None:
                try:
                    blob = shape.image.blob
                    ext = shape.image.ext or "png"
                    img_dir.mkdir(parents=True, exist_ok=True)
                    fn = img_dir / ("slide%02d-%02d.%s" % (si, len(images) + 1, ext))
                    fn.write_bytes(blob)
                    images.append(str(fn))
                    img_count += 1
                    continue
                except Exception:
                    pass
            # 표
            if shape.has_table:
                tables.append(extract_table(shape))
                continue
            # 텍스트
            lines = shape_text_lines(shape)
            if lines:
                if not title:
                    title = lines[0]["text"]
                    lines = lines[1:]
                body.extend(lines)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_out.append({
            "n": si,
            "title": title,
            "body": body,
            "tables": tables,
            "images": images,
            "notes": notes,
        })

    # ── 사람이 읽는 요약 출력 ─────────────────────────────────────────
    total = len(slides_out)
    print("# %s — %d 슬라이드, 이미지 %d개\n" % (src.name, total, img_count))
    for s in slides_out:
        print("=" * 64)
        print("[슬라이드 %02d] %s" % (s["n"], s["title"] or "(제목 없음)"))
        if s["body"]:
            print("  본문:")
            for ln in s["body"]:
                print("    " + "  " * ln["level"] + "- " + ln["text"])
        for ti, tbl in enumerate(s["tables"], 1):
            print("  표 %d (%d행):" % (ti, len(tbl)))
            for row in tbl:
                print("    | " + " | ".join(row) + " |")
        if s["images"]:
            print("  이미지: " + ", ".join(s["images"]))
        if s["notes"]:
            print("  노트(→ 설명란 재료):")
            for ln in s["notes"].splitlines():
                if ln.strip():
                    print("    " + ln.strip())
        print()

    if args.json:
        Path(args.json).write_text(
            json.dumps(slides_out, ensure_ascii=False, indent=2), encoding="utf-8")
        print("JSON 저장: %s" % args.json)


if __name__ == "__main__":
    main()
